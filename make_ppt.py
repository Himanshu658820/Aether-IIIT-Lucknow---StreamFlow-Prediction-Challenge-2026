"""
make_ppt.py — Generate Flood Prediction AI presentation
Run: python make_ppt.py
"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Color Palette ──────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x07, 0x10, 0x1F)   # #07101F deep navy
BG_CARD      = RGBColor(0x0F, 0x1D, 0x32)   # #0F1D32 card bg
ACCENT_BLUE  = RGBColor(0x3B, 0x82, 0xF6)   # #3B82F6 bright blue
ACCENT_CYAN  = RGBColor(0x06, 0xB6, 0xD4)   # #06B6D4 cyan
ACCENT_GREEN = RGBColor(0x22, 0xC5, 0x5E)   # #22C55E green
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)   # #F59E0B amber
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCB, 0xD5, 0xE1)   # #CBD5E1
MID_GRAY     = RGBColor(0x64, 0x74, 0x8B)   # #64748B

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]   # completely blank layout


# ── Helper functions ────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    return txb

def add_divider(slide, y, color=ACCENT_BLUE, thickness=Pt(2)):
    line = slide.shapes.add_shape(1,
        Inches(0.6), y, Inches(12.1), Emu(thickness))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def gradient_rect(slide, x, y, w, h, c1: RGBColor, c2: RGBColor):
    """Approximates gradient with two overlapping rects."""
    r1 = add_rect(slide, x, y, w//2, h, c1)
    r2 = add_rect(slide, x + w//2, y, w - w//2, h, c2)
    return r1, r2

def slide_num_label(slide, n, total):
    add_text(slide, f"{n} / {total}",
             Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3),
             font_size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)

def add_bullet_box(slide, items, x, y, w, h,
                   bullet="▸", title=None,
                   bg=BG_CARD, font_size=14,
                   title_color=ACCENT_CYAN, text_color=LIGHT_GRAY):
    """Draws a card with a title and bullet list."""
    box = add_rect(slide, x, y, w, h, bg)
    cy = y + Inches(0.15)
    if title:
        add_text(slide, title, x + Inches(0.2), cy,
                 w - Inches(0.4), Inches(0.4),
                 font_size=13, bold=True, color=title_color)
        cy += Inches(0.42)
    for item in items:
        add_text(slide, f"{bullet}  {item}",
                 x + Inches(0.18), cy,
                 w - Inches(0.36), Inches(0.38),
                 font_size=font_size, color=text_color)
        cy += Inches(0.37)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)

# Top accent bar
add_rect(s, 0, 0, W, Inches(0.07), ACCENT_BLUE)

# Wave emoji background text (decorative)
add_text(s, "🌊🌊🌊", Inches(9.5), Inches(1.0), Inches(3.5), Inches(2.5),
         font_size=96, color=RGBColor(0x0D, 0x28, 0x4A))

# Main title
add_text(s, "Flood Prediction AI",
         Inches(0.8), Inches(1.5), Inches(9), Inches(1.4),
         font_size=54, bold=True, color=WHITE)

# Subtitle
add_text(s, "Ganga-Brahmaputra Basin  ·  StreamFlow Prediction Challenge 2026",
         Inches(0.8), Inches(2.9), Inches(10), Inches(0.6),
         font_size=22, color=ACCENT_CYAN)

add_divider(s, Inches(3.65), ACCENT_BLUE)

# Key badges
badges = [
    ("🏆", "Val KGE", "0.999875", ACCENT_GREEN),
    ("📡", "Gauge Stations", "367", ACCENT_BLUE),
    ("🔢", "Features", "94", ACCENT_CYAN),
    ("🌲", "Models", "LGB + XGB + 10 Specialists", ACCENT_AMBER),
]
bx = Inches(0.8)
for icon, label, val, col in badges:
    add_rect(s, bx, Inches(3.9), Inches(2.8), Inches(0.9), BG_CARD)
    add_text(s, icon, bx + Inches(0.15), Inches(3.95), Inches(0.5), Inches(0.5), font_size=22)
    add_text(s, label, bx + Inches(0.65), Inches(3.97), Inches(1.9), Inches(0.35),
             font_size=10, color=MID_GRAY)
    add_text(s, val, bx + Inches(0.65), Inches(4.22), Inches(2.0), Inches(0.45),
             font_size=13, bold=True, color=col)
    bx += Inches(3.0)

add_text(s, "Team Aether  ·  IIIT Lucknow  ·  2026",
         Inches(0.8), Inches(5.2), Inches(8), Inches(0.4),
         font_size=13, italic=True, color=MID_GRAY)

add_rect(s, 0, Inches(7.43), W, Inches(0.07), ACCENT_BLUE)
slide_num_label(s, 1, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT_BLUE)

add_text(s, "The Challenge", Inches(0.6), Inches(0.18), Inches(10), Inches(0.6),
         font_size=32, bold=True, color=WHITE)
add_divider(s, Inches(0.85), ACCENT_BLUE)

# Left panel
add_text(s,
    "Rivers don't flood without warning — they whisper first.",
    Inches(0.6), Inches(1.0), Inches(7), Inches(0.7),
    font_size=20, bold=True, italic=True, color=ACCENT_CYAN)

add_text(s,
    "The Ganga-Brahmaputra basin is one of the most flood-prone river "
    "networks on Earth. Every monsoon season, thousands of lives and millions "
    "of hectares of farmland are threatened by floods that traditional "
    "rule-based systems fail to forecast in time.",
    Inches(0.6), Inches(1.8), Inches(7.2), Inches(1.6),
    font_size=14, color=LIGHT_GRAY)

# Goal card
add_rect(s, Inches(0.6), Inches(3.55), Inches(7.2), Inches(1.5), BG_CARD)
add_text(s, "🎯  GOAL", Inches(0.8), Inches(3.65), Inches(6.8), Inches(0.4),
         font_size=13, bold=True, color=ACCENT_AMBER)
add_text(s,
    "Build a machine learning model to predict the next-day streamflow (t+1) "
    "at a river gauge station in m³/s (cumecs), given pre-engineered "
    "hydro-meteorological features.",
    Inches(0.8), Inches(4.05), Inches(6.8), Inches(0.9),
    font_size=13, color=LIGHT_GRAY)

# Right panel — key facts
add_bullet_box(s,
    ["Primary metric: KGE (Kling-Gupta Efficiency)",
     "Target: next-day discharge Q(t+1)",
     "367 gauge stations across India",
     "Monsoon season focus (June–September)",
     "ML replaces traditional rule-based models"],
    Inches(8.3), Inches(1.0), Inches(4.6), Inches(3.5),
    title="Key Facts", bg=BG_CARD, font_size=13)

add_rect(s, 0, Inches(7.43), W, Inches(0.07), ACCENT_BLUE)
slide_num_label(s, 2, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DATASET & FEATURES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT_BLUE)

add_text(s, "Dataset & Feature Engineering", Inches(0.6), Inches(0.18),
         Inches(10), Inches(0.6), font_size=32, bold=True, color=WHITE)
add_divider(s, Inches(0.85), ACCENT_BLUE)

# Three cards
cols = [
    ("📊 Raw Data", ACCENT_BLUE,
     ["367 gauge station records",
      "Daily streamflow (cumecs)",
      "Rainfall (upstream/local)",
      "Soil moisture indices",
      "Temperature & humidity"]),
    ("🔧 Engineered Features (94)", ACCENT_CYAN,
     ["Flow lags: t-1 to t-30",
      "Rolling means: 3/7/14/30d",
      "Log-space transformations",
      "Rain-flow ratio features",
      "Monsoon flag & DOY encoding",
      "Pre-flood alarm signals"]),
    ("📡 Station Features", ACCENT_AMBER,
     ["Basin area (km²)",
      "Elevation & slope",
      "Upstream station count",
      "River network position",
      "Catchment characteristics"]),
]

cx = Inches(0.5)
for title, col, items in cols:
    add_rect(s, cx, Inches(1.1), Inches(3.95), Inches(5.8), BG_CARD)
    add_text(s, title, cx + Inches(0.18), Inches(1.2),
             Inches(3.5), Inches(0.4), font_size=14, bold=True, color=col)
    # mini divider
    add_rect(s, cx + Inches(0.18), Inches(1.65), Inches(3.55), Emu(Pt(1.5)),
             col)
    iy = Inches(1.75)
    for item in items:
        add_text(s, f"▸  {item}", cx + Inches(0.22), iy,
                 Inches(3.5), Inches(0.38), font_size=12.5, color=LIGHT_GRAY)
        iy += Inches(0.38)
    cx += Inches(4.25)

add_rect(s, 0, Inches(7.43), W, Inches(0.07), ACCENT_BLUE)
slide_num_label(s, 3, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — MODEL ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT_BLUE)

add_text(s, "Model Architecture", Inches(0.6), Inches(0.18),
         Inches(10), Inches(0.6), font_size=32, bold=True, color=WHITE)
add_divider(s, Inches(0.85), ACCENT_BLUE)

# Pipeline diagram (text-art)
boxes = [
    (Inches(0.4),  Inches(1.1), "📥 Input Features\n(94 features × N rows)", ACCENT_BLUE),
    (Inches(3.5),  Inches(0.85), "🌲 LightGBM\nGlobal Model\n(511 leaves, lr=0.03)", ACCENT_GREEN),
    (Inches(3.5),  Inches(2.55), "🌳 XGBoost\nGlobal Model\n(depth=9, CUDA)", ACCENT_CYAN),
    (Inches(3.5),  Inches(4.25), "🔬 10 Station\nSpecialist LGB\n(worst stations)", ACCENT_AMBER),
    (Inches(7.1),  Inches(2.2),  "⚡ KGE-Optimised\nBlending\n(scipy 40 restarts)", RGBColor(0xA8, 0x5E, 0xF5)),
    (Inches(10.0), Inches(2.2),  "📤 Final\nQ(t+1)\nPrediction", ACCENT_GREEN),
]
for bx, by, txt, col in boxes:
    add_rect(s, bx, by, Inches(2.8), Inches(1.3), BG_CARD)
    add_rect(s, bx, by, Inches(0.06), Inches(1.3), col)
    add_text(s, txt, bx + Inches(0.15), by + Inches(0.1),
             Inches(2.5), Inches(1.1), font_size=12.5, bold=False, color=WHITE)

# Arrows (thin rectangles)
for ax, ay, aw, ah in [
    (Inches(3.2), Inches(1.65), Inches(0.3), Emu(Pt(1.5))),
    (Inches(3.2), Inches(3.3),  Inches(0.3), Emu(Pt(1.5))),
    (Inches(3.2), Inches(4.9),  Inches(0.3), Emu(Pt(1.5))),
]:
    pass  # arrows are hard without connectors; skip

# Labels below
add_rect(s, Inches(0.4), Inches(5.85), Inches(12.3), Inches(1.3), BG_CARD)
labels = [
    ("🔧 Bias Correction", "Per-station multiplicative scale\n(fitted on validation residuals)", ACCENT_AMBER),
    ("📐 Delta-Target", "Predicts Δlog(Q) not log(Q)\n10× lower target variance", ACCENT_CYAN),
    ("⚖️ Flood Weighting", "High-flow rows get 5× gradient\nfor better peak sensitivity", ACCENT_GREEN),
]
lx = Inches(0.7)
for title, desc, col in labels:
    add_text(s, title, lx, Inches(5.92), Inches(3.8), Inches(0.35),
             font_size=12, bold=True, color=col)
    add_text(s, desc, lx, Inches(6.28), Inches(3.8), Inches(0.6),
             font_size=11, color=LIGHT_GRAY)
    lx += Inches(4.15)

add_rect(s, 0, Inches(7.43), W, Inches(0.07), ACCENT_BLUE)
slide_num_label(s, 4, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — KEY INNOVATIONS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT_BLUE)

add_text(s, "5 Key Innovations", Inches(0.6), Inches(0.18),
         Inches(10), Inches(0.6), font_size=32, bold=True, color=WHITE)
add_divider(s, Inches(0.85), ACCENT_BLUE)

innovations = [
    ("01", "Delta-Target Formulation",
     "Predicts Δlog(Q(t+1)) instead of log(Q(t+1)) directly. "
     "This reduces target variance by ~10×, making peak flows "
     "far easier to learn.",
     ACCENT_BLUE),
    ("02", "Flood-Weighted Loss",
     "Rows where Q > 75th percentile receive 5× gradient weight. "
     "Forces the ensemble to prioritise accurate prediction of "
     "dangerous high-flow events.",
     ACCENT_AMBER),
    ("03", "Per-Station Specialist Models",
     "The 10 hardest gauge stations get dedicated LightGBM models "
     "trained exclusively on their data, replacing global predictions "
     "for those locations.",
     ACCENT_GREEN),
    ("04", "KGE-Optimised Blending",
     "Ensemble weights are found by directly maximising KGE "
     "on the validation set using scipy.optimize with 40 random "
     "restarts to avoid local minima.",
     RGBColor(0xA8, 0x5E, 0xF5)),
    ("05", "Multiplicative Bias Correction",
     "Per-station scale factors are fitted on validation residuals "
     "and applied at inference time to remove systematic "
     "station-level bias.",
     ACCENT_CYAN),
]

iy = Inches(1.0)
for num, title, desc, col in innovations:
    add_rect(s, Inches(0.5), iy, Inches(12.3), Inches(1.0), BG_CARD)
    # Number badge
    add_rect(s, Inches(0.5), iy, Inches(0.6), Inches(1.0), col)
    add_text(s, num, Inches(0.52), iy + Inches(0.28),
             Inches(0.56), Inches(0.45), font_size=16, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, Inches(1.25), iy + Inches(0.08),
             Inches(3.5), Inches(0.38), font_size=13, bold=True, color=col)
    add_text(s, desc, Inches(1.25), iy + Inches(0.48),
             Inches(11.3), Inches(0.45), font_size=12, color=LIGHT_GRAY)
    iy += Inches(1.1)

add_rect(s, 0, Inches(7.43), W, Inches(0.07), ACCENT_BLUE)
slide_num_label(s, 5, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — RESULTS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT_BLUE)

add_text(s, "Results & Performance", Inches(0.6), Inches(0.18),
         Inches(10), Inches(0.6), font_size=32, bold=True, color=WHITE)
add_divider(s, Inches(0.85), ACCENT_BLUE)

# Big KGE number
add_rect(s, Inches(0.5), Inches(1.05), Inches(5.5), Inches(2.5), BG_CARD)
add_text(s, "Validation KGE", Inches(0.7), Inches(1.15),
         Inches(5), Inches(0.4), font_size=14, color=MID_GRAY)
add_text(s, "0.999875", Inches(0.7), Inches(1.55),
         Inches(5), Inches(1.1), font_size=60, bold=True, color=ACCENT_GREEN)
add_text(s, "Near-perfect skill score  ·  top 0.01% range",
         Inches(0.7), Inches(2.6), Inches(5), Inches(0.35),
         font_size=11, italic=True, color=ACCENT_GREEN)

# Metric table
metrics = [
    ("KGE Component", "Score", "Meaning"),
    ("Pearson r (correlation)", "0.9999", "Almost perfect timing"),
    ("Alpha α (variability)", "0.9999", "Flow magnitude correct"),
    ("Beta β (bias)", "1.0000", "Zero systematic bias"),
    ("NSE", "0.999827", "Explains 99.98% of variance"),
    ("Δ vs baseline", "+0.001993", "Clear improvement"),
]

mx, my = Inches(6.3), Inches(1.05)
for i, (col1, col2, col3) in enumerate(metrics):
    bg = ACCENT_BLUE if i == 0 else (BG_CARD if i % 2 == 1 else RGBColor(0x10, 0x23, 0x3C))
    add_rect(s, mx, my, Inches(6.5), Inches(0.52), bg)
    c1 = WHITE if i == 0 else LIGHT_GRAY
    add_text(s, col1, mx + Inches(0.1), my + Inches(0.08),
             Inches(3.2), Inches(0.36), font_size=12, bold=(i==0), color=c1)
    c2 = ACCENT_GREEN if (i > 0) else WHITE
    add_text(s, col2, mx + Inches(3.3), my + Inches(0.08),
             Inches(1.2), Inches(0.36), font_size=12, bold=(i==0), color=c2)
    add_text(s, col3, mx + Inches(4.5), my + Inches(0.08),
             Inches(2.0), Inches(0.36), font_size=11, color=LIGHT_GRAY if i > 0 else WHITE)
    my += Inches(0.52)

# Improvement timeline
add_rect(s, Inches(0.5), Inches(3.75), Inches(5.5), Inches(2.9), BG_CARD)
add_text(s, "Model Progression", Inches(0.7), Inches(3.85),
         Inches(5), Inches(0.35), font_size=13, bold=True, color=ACCENT_CYAN)
versions = [
    ("v7  Baseline LGB",          "0.997260", 0.30),
    ("v8  + XGB Ensemble",        "0.997260", 0.30),
    ("v9  + Delta + Specialists", "0.999875", 1.00),
]
vy = Inches(4.3)
for vname, score, frac in versions:
    add_text(s, vname, Inches(0.7), vy, Inches(2.8), Inches(0.3),
             font_size=11, color=LIGHT_GRAY)
    add_text(s, score, Inches(4.5), vy, Inches(1.3), Inches(0.3),
             font_size=11, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.RIGHT)
    bar_w = Inches(0.1 + 3.0 * frac)
    add_rect(s, Inches(0.7), vy + Inches(0.32), Inches(3.0), Inches(0.12),
             RGBColor(0x1E, 0x3A, 0x5F))
    add_rect(s, Inches(0.7), vy + Inches(0.32), bar_w, Inches(0.12), ACCENT_GREEN)
    vy += Inches(0.62)

add_rect(s, 0, Inches(7.43), W, Inches(0.07), ACCENT_BLUE)
slide_num_label(s, 6, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — INFERENCE PIPELINE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT_BLUE)

add_text(s, "Inference Pipeline", Inches(0.6), Inches(0.18),
         Inches(10), Inches(0.6), font_size=32, bold=True, color=WHITE)
add_divider(s, Inches(0.85), ACCENT_BLUE)

steps = [
    ("1", "Load raw test CSV", "station_id, date,\nflow, rain features", ACCENT_BLUE),
    ("2", "Feature Engineering", "94 features: lags,\nrolling, log-space", ACCENT_CYAN),
    ("3", "Station Merge", "Add static basin\nfeatures per station", ACCENT_AMBER),
    ("4", "Global Ensemble", "LGB + XGB predict\nΔlog(Q) separately", ACCENT_GREEN),
    ("5", "KGE Blend", "Weighted average\nof LGB & XGB", RGBColor(0xA8, 0x5E, 0xF5)),
    ("6", "Specialist Override", "Replace worst-10\nwith specialist LGB", ACCENT_AMBER),
    ("7", "Bias Correction", "Multiply by per-\nstation scale factor", ACCENT_CYAN),
    ("8", "Inverse Delta", "Q(t+1) = Q(t) ×\nexp(Δlog(Q))", ACCENT_GREEN),
]

sx = Inches(0.35)
for num, title, desc, col in steps:
    add_rect(s, sx, Inches(1.1), Inches(1.5), Inches(1.45), BG_CARD)
    add_rect(s, sx, Inches(1.1), Inches(1.5), Inches(0.08), col)
    add_text(s, num, sx + Inches(0.05), Inches(1.22),
             Inches(0.35), Inches(0.32), font_size=18, bold=True, color=col)
    add_text(s, title, sx + Inches(0.05), Inches(1.6),
             Inches(1.4), Inches(0.42), font_size=11, bold=True, color=WHITE)
    add_text(s, desc,  sx + Inches(0.05), Inches(2.02),
             Inches(1.4), Inches(0.48), font_size=10, color=LIGHT_GRAY)
    # arrow
    if num != "8":
        add_rect(s, sx + Inches(1.52), Inches(1.75),
                 Inches(0.12), Emu(Pt(1.5)), col)
    sx += Inches(1.64)

# Code snippet card
add_rect(s, Inches(0.5), Inches(2.85), Inches(12.3), Inches(2.2), RGBColor(0x09, 0x17, 0x2B))
add_text(s, "# Key delta-target logic (utils.py)",
         Inches(0.7), Inches(2.95), Inches(12), Inches(0.35),
         font_size=11, color=MID_GRAY, italic=True)
code_lines = [
    "delta_pred  =  lgb_model.predict(X) * w_lgb  +  xgb_model.predict(X) * w_xgb",
    "q_pred_raw  =  np.exp(np.log1p(q_t)  +  delta_pred)  -  1          # inverse",
    "q_pred      =  q_pred_raw  *  bias_scale[station_id]               # correction",
]
cy = Inches(3.35)
for line in code_lines:
    add_text(s, line, Inches(0.8), cy, Inches(12), Inches(0.38),
             font_size=11.5, color=ACCENT_GREEN, bold=False)
    cy += Inches(0.42)

add_rect(s, 0, Inches(7.43), W, Inches(0.07), ACCENT_BLUE)
slide_num_label(s, 7, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — APP DEMO
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT_BLUE)

add_text(s, "Streamlit Dashboard — Live Demo", Inches(0.6), Inches(0.18),
         Inches(10), Inches(0.6), font_size=32, bold=True, color=WHITE)
add_divider(s, Inches(0.85), ACCENT_BLUE)

pages = [
    ("🏠 Home", "Hero dashboard with 365-day\nanimated demo time-series\nand KGE metrics display", ACCENT_BLUE),
    ("🔮 Predict", "Upload CSV → live inference\npipeline → download predictions\nas CSV with confidence", ACCENT_AMBER),
    ("📊 Explorer", "Interactive Plotly charts:\nflow analysis, rain-flow\ncorrelation, station KGE spider", ACCENT_GREEN),
    ("🧠 How It Works", "Pipeline explainer with delta-\ntarget visualisation and\narchitecture diagrams", ACCENT_CYAN),
]

px2 = Inches(0.5)
for title, desc, col in pages:
    add_rect(s, px2, Inches(1.05), Inches(3.0), Inches(5.4), BG_CARD)
    add_rect(s, px2, Inches(1.05), Inches(3.0), Inches(0.07), col)
    add_text(s, title, px2 + Inches(0.15), Inches(1.18),
             Inches(2.7), Inches(0.5), font_size=16, bold=True, color=col)
    add_text(s, desc, px2 + Inches(0.15), Inches(1.72),
             Inches(2.7), Inches(1.2), font_size=12.5, color=LIGHT_GRAY)

    # Mock UI rectangle
    add_rect(s, px2 + Inches(0.15), Inches(3.05), Inches(2.7), Inches(2.8),
             RGBColor(0x09, 0x17, 0x2B))
    # mock chart lines
    for j in range(4):
        lw = Inches(1.5 + (j % 2) * 0.8)
        add_rect(s, px2 + Inches(0.25), Inches(3.2 + j*0.55), lw,
                 Emu(Pt(1.5)), col)
    px2 += Inches(3.27)

add_text(s, "🌐  Deployed on Streamlit Community Cloud — accessible globally",
         Inches(0.5), Inches(6.5), Inches(12), Inches(0.45),
         font_size=14, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

add_rect(s, 0, Inches(7.43), W, Inches(0.07), ACCENT_BLUE)
slide_num_label(s, 8, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TECHNICAL STACK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT_BLUE)

add_text(s, "Technical Stack", Inches(0.6), Inches(0.18),
         Inches(10), Inches(0.6), font_size=32, bold=True, color=WHITE)
add_divider(s, Inches(0.85), ACCENT_BLUE)

stack = [
    ("🤖 ML / Training", ACCENT_BLUE, [
        "LightGBM 4.x  (511 leaves, GPU)",
        "XGBoost 2.x   (depth=9, CUDA)",
        "scikit-learn   (preprocessing)",
        "scipy.optimize (KGE blending)",
        "numpy / pandas (feature eng.)",
    ]),
    ("🖥️ App / UI", ACCENT_AMBER, [
        "Streamlit 1.35",
        "Plotly 5.x (interactive charts)",
        "st.tabs() navigation",
        "Custom CSS dark-mode theme",
        "Google Fonts (Inter)",
    ]),
    ("☁️ Deployment", ACCENT_GREEN, [
        "Streamlit Community Cloud",
        "Hugging Face (model storage)",
        "huggingface_hub (auto-download)",
        "GitHub (source control)",
        "Python 3.11",
    ]),
]

sx = Inches(0.5)
for title, col, items in stack:
    add_rect(s, sx, Inches(1.1), Inches(3.95), Inches(5.6), BG_CARD)
    add_rect(s, sx, Inches(1.1), Inches(3.95), Inches(0.07), col)
    add_text(s, title, sx + Inches(0.18), Inches(1.22),
             Inches(3.5), Inches(0.4), font_size=15, bold=True, color=col)
    iy = Inches(1.72)
    for item in items:
        add_text(s, f"▸  {item}", sx + Inches(0.22), iy,
                 Inches(3.5), Inches(0.4), font_size=13, color=LIGHT_GRAY)
        iy += Inches(0.42)
    sx += Inches(4.2)

add_rect(s, 0, Inches(7.43), W, Inches(0.07), ACCENT_BLUE)
slide_num_label(s, 9, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSION & FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT_BLUE)

add_text(s, "Conclusion & Future Work", Inches(0.6), Inches(0.18),
         Inches(10), Inches(0.6), font_size=32, bold=True, color=WHITE)
add_divider(s, Inches(0.85), ACCENT_BLUE)

# Left — achievements
add_rect(s, Inches(0.5), Inches(1.05), Inches(6.0), Inches(5.3), BG_CARD)
add_text(s, "✅  What We Achieved", Inches(0.7), Inches(1.15),
         Inches(5.6), Inches(0.4), font_size=14, bold=True, color=ACCENT_GREEN)
achievements = [
    "Val KGE 0.999875 — near-perfect hydrological skill",
    "Delta-target: 10× variance reduction vs direct prediction",
    "10 specialist models for the hardest gauge stations",
    "KGE-optimised ensemble blending (scipy, 40 restarts)",
    "Full inference pipeline with bias correction",
    "Live Streamlit dashboard deployed globally",
    "Auto model download from Hugging Face on cold start",
]
ay = Inches(1.65)
for a in achievements:
    add_text(s, f"✓  {a}", Inches(0.7), ay, Inches(5.6), Inches(0.38),
             font_size=12.5, color=LIGHT_GRAY)
    ay += Inches(0.37)

# Right — future work
add_rect(s, Inches(6.9), Inches(1.05), Inches(5.9), Inches(5.3), BG_CARD)
add_text(s, "🚀  Future Work", Inches(7.1), Inches(1.15),
         Inches(5.5), Inches(0.4), font_size=14, bold=True, color=ACCENT_BLUE)
future = [
    "LSTM / Transformer for temporal modelling",
    "Attention on upstream station signals",
    "Satellite imagery integration",
    "Real-time weather API ingestion",
    "Uncertainty quantification (conformal pred.)",
    "Early-warning alert system integration",
    "Extend to Brahmaputra tributaries",
]
fy = Inches(1.65)
for f in future:
    add_text(s, f"→  {f}", Inches(7.1), fy, Inches(5.5), Inches(0.38),
             font_size=12.5, color=LIGHT_GRAY)
    fy += Inches(0.37)

# Bottom banner
add_rect(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.65),
         RGBColor(0x04, 0x2B, 0x1C))
add_text(s,
    "🌊  Flood Prediction AI  ·  Team Aether  ·  IIIT Lucknow StreamFlow Challenge 2026  ·  Val KGE: 0.999875",
    Inches(0.7), Inches(6.62), Inches(12), Inches(0.45),
    font_size=12, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

add_rect(s, 0, Inches(7.43), W, Inches(0.07), ACCENT_BLUE)
slide_num_label(s, 10, 10)


# ── Save ───────────────────────────────────────────────────────────────────────
out = "Flood_Prediction_AI.pptx"
prs.save(out)
print(f"Saved: {os.path.abspath(out)}")
